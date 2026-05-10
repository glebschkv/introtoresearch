#!/usr/bin/env python3
"""Extract text from all PDFs, PPTXs, DOCXs, and CSVs in the repo into one .txt file."""
import os
import sys
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

ROOT = Path(__file__).parent
OUTPUT = ROOT / "all_content.txt"

SEP = "=" * 80
SUB = "-" * 80


def header(title: str) -> str:
    return f"\n\n{SEP}\nFILE: {title}\n{SEP}\n\n"


def extract_pdf(path: Path) -> str:
    out = []
    reader = PdfReader(str(path))
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception as e:
            text = f"[error extracting page {i}: {e}]"
        out.append(f"{SUB}\n[Page {i}]\n{SUB}\n{text}")
    return "\n".join(out)


def extract_docx(path: Path) -> str:
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for t_idx, table in enumerate(doc.tables, 1):
        parts.append(f"\n[Table {t_idx}]")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"{SUB}\n[Slide {i}]\n{SUB}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        out.append(text)
            elif shape.shape_type == 19 or hasattr(shape, "text") and shape.has_text_frame is False:
                pass
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    out.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"[Speaker Notes]\n{notes}")
    return "\n".join(out)


def extract_csv(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".csv": extract_csv,
}


def main():
    files = sorted(
        p for p in ROOT.iterdir()
        if p.is_file() and p.suffix.lower() in EXTRACTORS
    )
    with OUTPUT.open("w", encoding="utf-8") as out:
        out.write("COMBINED COURSE MATERIALS — Introduction to Psychological Research\n")
        out.write(f"Source repo: introtoresearch\n")
        out.write(f"Total files: {len(files)}\n")
        out.write("\nFiles included:\n")
        for p in files:
            out.write(f"  - {p.name}\n")
        for p in files:
            ext = p.suffix.lower()
            print(f"Extracting {p.name}...", file=sys.stderr)
            try:
                text = EXTRACTORS[ext](p)
            except Exception as e:
                text = f"[FAILED TO EXTRACT: {e}]"
            out.write(header(p.name))
            out.write(text)
            out.write("\n")
    size = OUTPUT.stat().st_size
    print(f"\nWrote {OUTPUT} ({size:,} bytes)", file=sys.stderr)


if __name__ == "__main__":
    main()
